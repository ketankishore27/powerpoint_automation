from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime
from dateutil.relativedelta import relativedelta
import calendar

# === CONFIGURATION ===
num_months = 5  # Change to 3–12
start_month = datetime(2025, 6, 1)


# Determine columns and rows based on number of months
n_cols = 4 if num_months >= 10 else 3
n_rows = (num_months + n_cols - 1) // n_cols
legend_top_offset_inch = 0.4 if n_rows <= 2 else 0.25
# === Events Data ===
events_by_type = {
    "CPI": ["2025-06-11", "2025-07-15", "2025-08-12", "2025-09-11", "2025-10-15", "2025-11-13", "2025-12-10"],
    "PPI": ["2025-06-12", "2025-07-16", "2025-08-14", "2025-09-10", "2025-10-16", "2025-11-14", "2025-12-11"],
    "HLD": ["2025-07-04", "2025-09-01", "2025-10-13", "2025-11-11", "2025-11-27", "2025-12-25"]
}
event_lookup = {date: label for label, dates in events_by_type.items() for date in dates}

# === Sizing and Fonts ===
slide_width = 10.0
slide_height = 7.5

calendar_width_in = (slide_width - 1.5) / n_cols
calendar_height_in = (slide_height - 2.0) / n_rows

cell_width = Inches(calendar_width_in / 7)
cell_height = Inches(calendar_height_in / 8)

base_font = min(cell_width.inches, cell_height.inches) * 20
title_font = int(base_font * 5)
month_font = int(base_font * 2)
day_font = int(base_font * 1.1)
legend_font = int(base_font * 1.5)
legend_box_size = Inches(0.2)

colors = {
    "PPI": RGBColor(217, 217, 217),
    "CPI": RGBColor(217, 217, 217),
    "HLD": RGBColor(255, 255, 224),
    "WEEKDAY_BG": RGBColor(0, 51, 102),
    "MONTH_TEXT": RGBColor(70, 130, 180),
    "DAY_BOX": RGBColor(255, 255, 255),
}

months = [(start_month + relativedelta(months=i)).strftime("%B") for i in range(num_months)]
years = [(start_month + relativedelta(months=i)).year for i in range(num_months)]

# === Create Presentation ===
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Title
title_shape = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.4))
title_shape.text = "Calendar"
title_shape.text_frame.paragraphs[0].font.size = Pt(title_font)
title_shape.text_frame.paragraphs[0].font.bold = True

# Margins and offsets
start_left = Inches(0.3)
start_top = Inches(0.9)
x_margin = Inches(0.2)
y_margin = Inches(0.3)

# === Draw Each Calendar Month ===
for idx, (month_name, year) in enumerate(zip(months, years)):
    month_index = datetime.strptime(month_name, "%B").month
    cal = calendar.Calendar(firstweekday=0)
    days = list(cal.itermonthdays(year, month_index))

    grid_col = idx % n_cols
    grid_row = idx // n_cols
    x_offset = start_left + grid_col * (cell_width * 7 + x_margin)
    y_offset = start_top + grid_row * (cell_height * 8 + y_margin)

    # Month label
    month_box = slide.shapes.add_textbox(x_offset, y_offset, cell_width * 7, Inches(0.3))
    tf = month_box.text_frame
    tf.text = f"{month_name} {year}"
    tf.paragraphs[0].font.size = Pt(month_font)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = colors["MONTH_TEXT"]
    y_offset += Inches(0.3)

    # Weekday headers
    weekdays = ["M", "T", "W", "Th", "F", "Sa", "Su"]
    for i, day in enumerate(weekdays):
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_offset + i * cell_width, y_offset, cell_width, cell_height)
        box.fill.solid()
        box.fill.fore_color.rgb = colors["WEEKDAY_BG"]
        box.line.color.rgb = RGBColor(0, 0, 0)
        box.text = day
        p = box.text_frame.paragraphs[0]
        p.font.size = Pt(day_font)
        p.font.color.rgb = RGBColor(255, 255, 255)

    y_offset += cell_height + Inches(0.02)

    # Day cells
    col, row = 0, 0
    for day in days:
        left = x_offset + col * cell_width
        top = y_offset + row * cell_height
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, cell_width, cell_height)
        box.line.color.rgb = RGBColor(100, 100, 100)
        box.fill.solid()
        if day == 0:
            box.fill.fore_color.rgb = colors["DAY_BOX"]
        else:
            date_str = f"{year}-{month_index:02d}-{day:02d}"
            label = event_lookup.get(date_str)
            box.fill.fore_color.rgb = colors[label] if label else colors["DAY_BOX"]
            tf = box.text_frame
            tf.text = f"{day}"
            tf.paragraphs[0].font.size = Pt(day_font)
            tf.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            if label:
                p2 = tf.add_paragraph()
                p2.text = label
                p2.font.size = Pt(day_font - 1)
                p2.font.color.rgb = RGBColor(0, 0, 0)
        col += 1
        if col >= 7:
            col = 0
            row += 1

# === Legend ===
legend_top = Inches(slide_height - legend_top_offset_inch)
legend_left = start_left
legend_gap = Inches(1.2)

for i, key in enumerate(["PPI", "CPI", "HLD"]):
    color_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, legend_left + i * legend_gap, legend_top, legend_box_size, legend_box_size)
    color_box.fill.solid()
    color_box.fill.fore_color.rgb = colors[key]
    color_box.line.color.rgb = RGBColor(0, 0, 0)

    label_box = slide.shapes.add_textbox(legend_left + i * legend_gap + Inches(0.25), legend_top, Inches(1), legend_box_size)
    tf = label_box.text_frame
    tf.text = key
    tf.paragraphs[0].font.size = Pt(legend_font)

# === Save Presentation ===
prs.save("calendar_adaptive_layout.pptx")