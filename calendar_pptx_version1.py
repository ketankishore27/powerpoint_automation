from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime
import calendar

# 🎨 Color mapping
colors = {
    "PPI": RGBColor(217, 217, 217),     # Light grey
    "CPI": RGBColor(217, 217, 217),     # Light grey
    "HLD": RGBColor(255, 255, 224),     # Light yellow
    "WEEKDAY_BG": RGBColor(0, 51, 102),  # Dark blue
    "MONTH_TEXT": RGBColor(70, 130, 180),  # Sky blue
    "DAY_BOX": RGBColor(255, 255, 255),   # White background for all day numbers
    "EMPTY_BOX": RGBColor(240, 240, 240)  # Light grey for empty days
}

# ✅ Input: Event dictionary by label
events_by_type = {
    "CPI": ["2025-06-11", "2025-07-15", "2025-08-12", "2025-09-11", "2025-10-15", "2025-11-13"],
    "PPI": ["2025-06-12", "2025-07-16", "2025-08-14", "2025-09-10", "2025-10-16", "2025-11-14"],
    "HLD": ["2025-07-04", "2025-09-01", "2025-10-13", "2025-11-11", "2025-11-27"]
}

# ➕ Convert to date-to-label mapping
event_lookup = {}
for label, date_list in events_by_type.items():
    for d in date_list:
        event_lookup[d] = label

# 📄 Setup presentation
prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

# 📌 Title
title_shape = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
title_shape.text = "Execution Calendar Window"
title_shape.text_frame.paragraphs[0].font.size = Pt(24)
title_shape.text_frame.paragraphs[0].font.bold = True

# 🔧 Calendar layout
cell_width = Inches(0.4)
cell_height = Inches(0.3)
start_left = Inches(0.4)
start_top = Inches(1.0)
x_margin = Inches(0.25)
y_margin = Inches(0.5)

months = ["June", "July", "August", "September", "October", "November"]

for idx, month_name in enumerate(months):
    month_index = datetime.strptime(month_name, "%B").month
    cal = calendar.Calendar(firstweekday=0)
    month_days = list(cal.itermonthdays(2025, month_index))

    grid_col = idx % 3
    grid_row = idx // 3
    x_offset = start_left + grid_col * (cell_width * 7 + x_margin)
    y_offset = start_top + grid_row * (cell_height * 8 + y_margin)

    # 🗓️ Month title
    month_textbox = slide.shapes.add_textbox(x_offset, y_offset, cell_width * 7, Inches(0.3))
    tf = month_textbox.text_frame
    tf.text = f"{month_name} 2025"
    p = tf.paragraphs[0]
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = colors["MONTH_TEXT"]
    y_offset += Inches(0.35)

    # 📆 Weekday headers
    weekdays = ["M", "T", "W", "Th", "F", "Sa", "Su"]
    for i, day in enumerate(weekdays):
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_offset + i * cell_width, y_offset, cell_width, cell_height)
        box.fill.solid()
        box.fill.fore_color.rgb = colors["WEEKDAY_BG"]
        box.line.color.rgb = RGBColor(0, 0, 0)
        box.text = day
        p = box.text_frame.paragraphs[0]
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(255, 255, 255)  # White text

    y_offset += cell_height + Inches(0.05)  # Space between day names and dates

    # 📅 Day boxes
    col = 0
    row = 0
    for day in month_days:
        left = x_offset + col * cell_width
        top = y_offset + row * cell_height
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, cell_width, cell_height)
        box.line.color.rgb = RGBColor(100, 100, 100)

        if day == 0:
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background for empty boxes
            box.line.color.rgb = RGBColor(100, 100, 100)
        else:
            date_key = f"2025-{month_index:02d}-{day:02d}"
            label = event_lookup.get(date_key, "")
            box.fill.solid()

            # 🎨 Color based on event type
            if label:
                box.fill.fore_color.rgb = colors[label]
            else:
                box.fill.fore_color.rgb = colors["DAY_BOX"]

            tf = box.text_frame
            tf.text = f"{day}"
            tf.paragraphs[0].font.size = Pt(8)
            tf.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

            if label:
                p2 = tf.add_paragraph()
                p2.text = label
                p2.font.size = Pt(7)
                p2.font.color.rgb = RGBColor(0, 0, 0)

        col += 1
        if col >= 7:
            col = 0
            row += 1

# 🧾 Legend
legend_top = start_top + 2 * (cell_height * 8 + y_margin) + Inches(0.2)  # Removed 4 inch padding
legend_left = start_left
legend_gap = Inches(1.2)

for i, key in enumerate(["PPI", "CPI", "HLD"]):
    color_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, legend_left + i * legend_gap, legend_top, Inches(0.25), Inches(0.25))
    color_box.fill.solid()
    color_box.fill.fore_color.rgb = colors[key]
    color_box.line.color.rgb = RGBColor(0, 0, 0)

    label_box = slide.shapes.add_textbox(legend_left + i * legend_gap + Inches(0.3), legend_top, Inches(1), Inches(0.25))
    tf = label_box.text_frame
    tf.text = key
    tf.paragraphs[0].font.size = Pt(9)

# ✅ Save Presentation
prs.save("calendar_final_with_legends.pptx")
print("✅ Saved as: calendar_final_with_legends.pptx")